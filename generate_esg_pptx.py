from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Color scheme
colors = {
    'dark_green': RGBColor(28, 89, 61),
    'light_green': RGBColor(76, 175, 80),
    'blue': RGBColor(0, 150, 200),
    'cream': RGBColor(240, 240, 235),
    'white': RGBColor(255, 255, 255)
}

# Create a PowerPoint presentation
prs = Presentation()

# Function to create slide
def create_slide(title_text, content_text, background_color):
    slide_layout = prs.slide_layouts[5]  # Title Slide Layout
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = background_color

    title = slide.shapes.title
    content = slide.placeholders[0]

    title.text = title_text
    content.text = content_text
    return slide

# Slide content
slides_content = [
    ("ArrowGreenTech ESG Report", "Generated on: 2026-03-18", colors['dark_green']),
    ("Mission Statement", "To promote sustainable growth through comprehensive ESG strategies.", colors['white']),
    ("Overview of ESG", "Environmental, Social, and Governance criteria are essential for sustainability.", colors['cream']),
    ("Importance of ESG Reporting", "Provides transparency and accountability to stakeholders.", colors['light_green']),
    ("Corporate Governance", "Framework for achieving a company's objectives.", colors['blue']),
    ("Environmental Impact", "Measures taken to reduce carbon footprint and promote green energy.", colors['dark_green']),
    ("Social Responsibility", "Community engagement and social initiatives for welfare.", colors['white']),
    ("Economic Value", "Balancing profit and purpose to generate long-term impact.", colors['cream']),
    ("Key Performance Indicators (KPIs)", "Metrics to evaluate the success of ESG initiatives.", colors['light_green']),
    ("Stakeholder Engagement", "Building trust with transparency in communication.", colors['blue']),
    ("Materiality Assessment Introduction", "Identifying what matters most to our stakeholders.", colors['dark_green']),
    ("Materiality Assessment Details", "In-depth analysis of key areas of concern.", colors['white']),
    ("ESG Challenges", "Addressing the roadblocks in our ESG journey.", colors['cream']),
    ("ESG Initiatives", "Programs in place to enhance ESG performance.", colors['light_green']),
    ("Future Goals", "Setting ambitious targets for growth and sustainability.", colors['blue']),
    ("Conclusion", "Commitment to continuous improvement and transparency.", colors['dark_green']),
]

# Create slides based on the content
for title, content, color in slides_content:
    create_slide(title, content, color)

# Save the presentation
prs.save("ArrowGreenTech_ESG_Report.pptx")
